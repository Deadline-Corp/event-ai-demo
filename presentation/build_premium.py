"""
Premium EVENT AI Investor Deck — python-pptx builder
Strategy: full-bleed PNG backgrounds (visual richness) + real text boxes on top (editable)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ── Paths ──────────────────────────────────────────────────────────────────────
BASE      = r"D:\event-ai\presentation"
BG_DIR    = os.path.join(BASE, "bg", "png")
LOGO      = r"D:\event-ai\assets\logo-dark.png"
OUT_PPTX  = os.path.join(BASE, "EVENT-AI-Investor-RU.pptx")

BG_COVER   = os.path.join(BG_DIR, "bg-cover__1920x1080.png")
BG_CONTENT = os.path.join(BG_DIR, "bg-content__1920x1080.png")
BG_ACCENT  = os.path.join(BG_DIR, "bg-accent__1920x1080.png")

# ── Slide dimensions: 16:9  13.333 × 7.5 inches ───────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

# ── Color palette ──────────────────────────────────────────────────────────────
C_WHITE   = RGBColor(0xEE, 0xF1, 0xF8)   # primary text
C_DIM     = RGBColor(0x9A, 0xA4, 0xB6)   # secondary text
C_ACCENT  = RGBColor(0x2C, 0x7B, 0xFF)   # blue accent
C_ACCENT2 = RGBColor(0x5A, 0xA0, 0xFF)   # lighter blue
C_DARK    = RGBColor(0x04, 0x06, 0x0B)   # background
C_CARD_BG = RGBColor(0x0C, 0x10, 0x1A)   # card fill

# ── Fonts ──────────────────────────────────────────────────────────────────────
F_HEAD  = "Space Grotesk"
F_BODY  = "Montserrat"
F_MONO  = "JetBrains Mono"


# ══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ══════════════════════════════════════════════════════════════════════════════

def add_bg(slide, png_path):
    """Insert a full-bleed PNG as slide background, sent to back."""
    pic = slide.shapes.add_picture(png_path, 0, 0, W, H)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def add_textbox(slide, text, left, top, width, height,
                font_name=F_BODY, font_size=18, bold=False, italic=False,
                color=C_WHITE, align=PP_ALIGN.LEFT,
                word_wrap=True, line_spacing=None):
    """Add a single-paragraph textbox."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name  = font_name
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    if italic:
        run.font.italic = italic
    return txBox


def add_multiline_textbox(slide, lines, left, top, width, height,
                          font_name=F_BODY, font_size=18, bold=False,
                          color=C_WHITE, align=PP_ALIGN.LEFT,
                          line_spacing=1.15, para_space_before=Pt(4)):
    """Add a textbox with multiple paragraphs (one per list item)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = para_space_before
        run = p.add_run()
        run.text = line
        run.font.name  = font_name
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height,
             fill_color=None, fill_alpha=None,
             line_color=None, line_width_pt=0.75):
    """Add a rectangle shape (card / divider)."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()

    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = Pt(line_width_pt)
    else:
        line.fill.background()

    return shape


def add_thin_line(slide, left, top, width, color=C_ACCENT, thickness_pt=1.0):
    """Horizontal thin accent line."""
    shape = slide.shapes.add_shape(1, left, top, width, Pt(thickness_pt))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def set_fill_transparency(shape, alpha_pct):
    """
    Patch a shape's solid fill to have transparency (0=opaque, 100=clear).
    alpha_pct: 0..100
    """
    solidFill = shape.fill._xPr.find(qn('a:solidFill'))
    if solidFill is None:
        return
    srgbClr = solidFill.find(qn('a:srgbClr'))
    if srgbClr is None:
        return
    # Remove existing alpha if present
    for a in srgbClr.findall(qn('a:alpha')):
        srgbClr.remove(a)
    alpha_el = etree.SubElement(srgbClr, qn('a:alpha'))
    # pptx alpha: 100% = 100000, 0% = 0
    val = int((100 - alpha_pct) * 1000)
    alpha_el.set('val', str(val))


def add_circle_badge(slide, cx, cy, r, number,
                     bg_color=C_ACCENT, text_color=C_WHITE):
    """Filled circle with a number — step indicators."""
    d = r * 2
    shape = slide.shapes.add_shape(9, cx - r, cy - r, d, d)  # 9 = oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.name  = F_MONO
    run.font.size  = Pt(14)
    run.font.bold  = True
    run.font.color.rgb = text_color
    return shape


def add_logo(slide, left, top, height):
    """Insert logo PNG at given position, maintain aspect ratio."""
    slide.shapes.add_picture(LOGO, left, top, height=height)


# ══════════════════════════════════════════════════════════════════════════════
# Slide builders
# ══════════════════════════════════════════════════════════════════════════════

MARGIN_L  = Inches(0.75)
MARGIN_R  = Inches(0.75)
MARGIN_T  = Inches(0.55)
SAFE_W    = W - MARGIN_L - MARGIN_R  # ~11.83 in

SMALL_LOGO_H = Inches(0.38)
SMALL_LOGO_T = Inches(0.22)
SMALL_LOGO_L = Inches(0.65)

def add_small_logo(slide):
    add_logo(slide, SMALL_LOGO_L, SMALL_LOGO_T, SMALL_LOGO_H)


# ─────────────────────────────────────────────────────────────────────────────
# S1 — Cover
# ─────────────────────────────────────────────────────────────────────────────
def build_s1(slide):
    add_bg(slide, BG_COVER)

    # Logo — large, centered, upper area
    logo_w = Inches(4.8)
    logo_l = (W - logo_w) / 2
    logo_t = Inches(1.5)
    slide.shapes.add_picture(LOGO, logo_l, logo_t, width=logo_w)

    # Thin accent line under logo
    add_thin_line(slide, (W - Inches(1.4)) / 2, Inches(2.95), Inches(1.4),
                  color=C_ACCENT, thickness_pt=2.0)

    # Kicker
    add_textbox(slide, "ИНВЕСТИЦИОННАЯ ПРЕЗЕНТАЦИЯ",
                MARGIN_L, Inches(3.18), SAFE_W, Inches(0.45),
                font_name=F_MONO, font_size=12, bold=False,
                color=C_ACCENT2, align=PP_ALIGN.CENTER)

    # Main subtitle — centered in a narrower box to force visual centering
    sub_w = Inches(9.0)
    sub_l = (W - sub_w) / 2
    add_textbox(slide,
                "Умный AI-помощник для организации\nмероприятий под ключ",
                sub_l, Inches(3.68), sub_w, Inches(1.5),
                font_name=F_HEAD, font_size=30, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, "Алматы  ·  Казахстан  ·  2026",
                MARGIN_L, Inches(6.6), SAFE_W, Inches(0.45),
                font_name=F_BODY, font_size=14, bold=False,
                color=C_DIM, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# S2 — Problem
# ─────────────────────────────────────────────────────────────────────────────
def build_s2(slide):
    add_bg(slide, BG_CONTENT)
    add_small_logo(slide)

    # Slide number badge
    add_textbox(slide, "01", W - Inches(1.1), MARGIN_T, Inches(0.8), Inches(0.4),
                font_name=F_MONO, font_size=11, color=C_DIM, align=PP_ALIGN.RIGHT)

    # Section kicker
    add_textbox(slide, "ПРОБЛЕМА", MARGIN_L, Inches(0.55), SAFE_W, Inches(0.35),
                font_name=F_MONO, font_size=11, color=C_ACCENT2)

    # Title
    add_textbox(slide, "Организация мероприятий — сложно и долго",
                MARGIN_L, Inches(0.95), SAFE_W, Inches(1.1),
                font_name=F_HEAD, font_size=32, bold=True,
                color=C_WHITE, word_wrap=True)

    # Thin line
    add_thin_line(slide, MARGIN_L, Inches(2.2), Inches(0.7))

    # Pain points — 5 cards in 2 rows (3 top + 2 bottom centered)
    pains = [
        ("01", "Не знаешь, с чего начать"),
        ("02", "Площадку, ведущих, артистов\nищешь по отдельности"),
        ("03", "Нет прозрачной сметы"),
        ("04", "Сложно сравнить\nподрядчиков"),
        ("05", "Время уходит на звонки\nи согласования"),
    ]

    card_w  = Inches(3.7)
    card_h  = Inches(1.65)
    gap     = Inches(0.18)
    row1_y  = Inches(2.42)
    row2_y  = row1_y + card_h + gap
    # Row 1: 3 cards
    row1_total = 3 * card_w + 2 * gap
    row1_x_start = (W - row1_total) / 2

    for i in range(3):
        cx = row1_x_start + i * (card_w + gap)
        _build_pain_card(slide, pains[i], cx, row1_y, card_w, card_h)

    # Row 2: 2 cards centered
    row2_total = 2 * card_w + gap
    row2_x_start = (W - row2_total) / 2
    for i in range(2):
        cx = row2_x_start + i * (card_w + gap)
        _build_pain_card(slide, pains[3 + i], cx, row2_y, card_w, card_h)


def _build_pain_card(slide, pain_tuple, left, top, width, height):
    num, text = pain_tuple
    # Card bg
    card = add_rect(slide, left, top, width, height,
                    fill_color=RGBColor(0x0A, 0x0F, 0x1C),
                    line_color=RGBColor(0x1A, 0x28, 0x4D), line_width_pt=0.75)
    try:
        set_fill_transparency(card, 0)  # fully opaque dark card
    except Exception:
        pass

    # Number
    add_textbox(slide, num,
                left + Inches(0.22), top + Inches(0.18),
                Inches(0.6), Inches(0.35),
                font_name=F_MONO, font_size=11, color=C_ACCENT2)

    # Text
    add_textbox(slide, text,
                left + Inches(0.22), top + Inches(0.6),
                width - Inches(0.44), height - Inches(0.75),
                font_name=F_HEAD, font_size=17, bold=True,
                color=C_WHITE, word_wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# S3 — Solution
# ─────────────────────────────────────────────────────────────────────────────
def build_s3(slide):
    add_bg(slide, BG_CONTENT)
    add_small_logo(slide)

    add_textbox(slide, "02", W - Inches(1.1), MARGIN_T, Inches(0.8), Inches(0.4),
                font_name=F_MONO, font_size=11, color=C_DIM, align=PP_ALIGN.RIGHT)

    add_textbox(slide, "РЕШЕНИЕ", MARGIN_L, Inches(0.55), SAFE_W, Inches(0.35),
                font_name=F_MONO, font_size=11, color=C_ACCENT2)

    add_textbox(slide, "EVENT AI собирает мероприятие за вас",
                MARGIN_L, Inches(0.95), SAFE_W, Inches(0.75),
                font_name=F_HEAD, font_size=36, bold=True, color=C_WHITE)

    add_textbox(slide, "Ответьте на пару вопросов в AI-чате — система соберёт:",
                MARGIN_L, Inches(1.75), SAFE_W, Inches(0.45),
                font_name=F_BODY, font_size=16, color=C_DIM)

    add_thin_line(slide, MARGIN_L, Inches(2.3), Inches(0.7))

    # No emoji — use bullet squares that render reliably in pptx
    features = [
        ("#", "Подборку\nподрядчиков"),
        ("$", "Смету"),
        ("%", "Сценарий"),
        ("@", "Тайминг"),
        ("&", "Пригласительные"),
        ("*", "Бронирование"),
    ]

    card_w = Inches(3.7)
    card_h = Inches(1.72)
    gap    = Inches(0.20)
    rows   = [Inches(2.5), Inches(2.5) + card_h + gap]
    cols   = [MARGIN_L,
              MARGIN_L + card_w + gap,
              MARGIN_L + 2 * (card_w + gap)]

    for idx, (icon, label) in enumerate(features):
        row = idx // 3
        col = idx % 3
        cx = cols[col]
        cy = rows[row]
        _build_feature_card(slide, label, cx, cy, card_w, card_h)


def _build_feature_card(slide, label, left, top, width, height):
    card = add_rect(slide, left, top, width, height,
                    fill_color=RGBColor(0x0A, 0x0F, 0x1C),
                    line_color=RGBColor(0x1A, 0x28, 0x4D), line_width_pt=0.75)

    # Accent dot
    dot = add_rect(slide, left + Inches(0.22), top + Inches(0.22),
                   Inches(0.08), Inches(0.08),
                   fill_color=C_ACCENT, line_color=None)

    # Label
    add_textbox(slide, label,
                left + Inches(0.22), top + Inches(0.42),
                width - Inches(0.44), height - Inches(0.55),
                font_name=F_HEAD, font_size=18, bold=True,
                color=C_WHITE, word_wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# S4 — How it works (steps)
# ─────────────────────────────────────────────────────────────────────────────
def build_s4(slide):
    add_bg(slide, BG_CONTENT)
    add_small_logo(slide)

    add_textbox(slide, "03", W - Inches(1.1), MARGIN_T, Inches(0.8), Inches(0.4),
                font_name=F_MONO, font_size=11, color=C_DIM, align=PP_ALIGN.RIGHT)

    add_textbox(slide, "КАК ЭТО РАБОТАЕТ", MARGIN_L, Inches(0.55), SAFE_W, Inches(0.35),
                font_name=F_MONO, font_size=11, color=C_ACCENT2)

    add_textbox(slide, "От идеи до брони — в одном приложении",
                MARGIN_L, Inches(0.95), SAFE_W, Inches(0.9),
                font_name=F_HEAD, font_size=38, bold=True, color=C_WHITE)

    add_thin_line(slide, MARGIN_L, Inches(2.0), Inches(0.7))

    steps = [
        ("1", "Выбираете тип\nмероприятия"),
        ("2", "AI уточняет\nдетали"),
        ("3", "Получаете\nподборку"),
        ("4", "Выбираете\nпроекты"),
        ("5", "Бронируете"),
    ]

    card_w  = Inches(2.28)
    card_h  = Inches(3.2)
    gap     = Inches(0.20)
    total_w = 5 * card_w + 4 * gap
    start_x = (W - total_w) / 2
    card_y  = Inches(2.3)

    for i, (num, label) in enumerate(steps):
        cx = start_x + i * (card_w + gap)
        _build_step_card(slide, num, label, cx, card_y, card_w, card_h, i)

    # Arrow connectors between cards — ">" text centered in gap
    for i in range(4):
        ax = start_x + (i + 1) * card_w + i * gap
        ay = card_y + card_h / 2 - Inches(0.18)
        add_textbox(slide, ">", ax, ay, gap, Inches(0.36),
                    font_name=F_MONO, font_size=16, bold=True,
                    color=C_ACCENT2, align=PP_ALIGN.CENTER)


def _build_step_card(slide, num, label, left, top, width, height, step_idx):
    # Slightly highlight first card
    bord = RGBColor(0x2C, 0x7B, 0xFF) if step_idx == 0 else RGBColor(0x1A, 0x28, 0x4D)
    card = add_rect(slide, left, top, width, height,
                    fill_color=RGBColor(0x0A, 0x0F, 0x1C),
                    line_color=bord, line_width_pt=1.0)

    # Circle number
    r = Inches(0.32)
    cx = left + width / 2
    cy = top + Inches(0.65)
    add_circle_badge(slide, cx, cy, r, num, bg_color=C_ACCENT)

    # Label
    add_textbox(slide, label,
                left + Inches(0.12), top + Inches(1.3),
                width - Inches(0.24), height - Inches(1.5),
                font_name=F_HEAD, font_size=17, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER, word_wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# S5 — Monetization
# ─────────────────────────────────────────────────────────────────────────────
def build_s5(slide):
    add_bg(slide, BG_ACCENT)
    add_small_logo(slide)

    add_textbox(slide, "04", W - Inches(1.1), MARGIN_T, Inches(0.8), Inches(0.4),
                font_name=F_MONO, font_size=11, color=C_DIM, align=PP_ALIGN.RIGHT)

    add_textbox(slide, "МОНЕТИЗАЦИЯ", MARGIN_L, Inches(0.55), SAFE_W, Inches(0.35),
                font_name=F_MONO, font_size=11, color=C_ACCENT2)

    add_textbox(slide, "Как зарабатывает платформа",
                MARGIN_L, Inches(0.95), SAFE_W, Inches(0.9),
                font_name=F_HEAD, font_size=38, bold=True, color=C_WHITE)

    add_thin_line(slide, MARGIN_L, Inches(2.0), Inches(0.7))

    # Two big metric cards side by side
    metric_card_w = Inches(4.5)
    metric_card_h = Inches(2.1)
    gap = Inches(0.35)
    metrics_total = 2 * metric_card_w + gap
    mx_start = (W - metrics_total) / 2 - Inches(0.3)
    my = Inches(2.25)

    _build_metric_card(slide, "7%", "Комиссия с каждого заказа",
                       mx_start, my, metric_card_w, metric_card_h)
    _build_metric_card(slide, "50 000 ₸", "Стоимость бронирования",
                       mx_start + metric_card_w + gap, my, metric_card_w, metric_card_h)

    # 3 revenue streams below
    streams = [
        "Продвижение партнёров",
        "Премиум-размещение",
        "B2B-пакеты для корпоратов",
    ]
    stream_card_w = Inches(3.7)
    stream_card_h = Inches(0.95)
    s_gap = Inches(0.22)
    total_sw = 3 * stream_card_w + 2 * s_gap
    sx_start = (W - total_sw) / 2
    sy = Inches(4.65)

    for i, txt in enumerate(streams):
        cx = sx_start + i * (stream_card_w + s_gap)
        card = add_rect(slide, cx, sy, stream_card_w, stream_card_h,
                        fill_color=RGBColor(0x0A, 0x0F, 0x1C),
                        line_color=RGBColor(0x1A, 0x28, 0x4D), line_width_pt=0.75)
        add_textbox(slide, txt,
                    cx + Inches(0.18), sy + Inches(0.2),
                    stream_card_w - Inches(0.36), stream_card_h - Inches(0.3),
                    font_name=F_HEAD, font_size=16, bold=False,
                    color=C_WHITE, align=PP_ALIGN.CENTER, word_wrap=True)

    # Flow diagram: Клиент → Платформа → Партнёр
    flow_y = Inches(5.85)
    flow_items = ["Клиент", "Платформа", "Партнёр"]
    flow_w = Inches(2.5)
    flow_h = Inches(0.65)
    arrow_w = Inches(0.8)
    total_fw = 3 * flow_w + 2 * arrow_w
    fx_start = (W - total_fw) / 2

    for i, label in enumerate(flow_items):
        x = fx_start + i * (flow_w + arrow_w)
        fill_c = C_ACCENT if i == 1 else RGBColor(0x0C, 0x12, 0x20)
        bord_c = C_ACCENT if i == 1 else RGBColor(0x1A, 0x28, 0x4D)
        card = add_rect(slide, x, flow_y, flow_w, flow_h,
                        fill_color=fill_c, line_color=bord_c, line_width_pt=1.0)
        text_c = C_WHITE
        add_textbox(slide, label, x, flow_y + Inches(0.14),
                    flow_w, Inches(0.4),
                    font_name=F_HEAD, font_size=16, bold=True,
                    color=text_c, align=PP_ALIGN.CENTER)

        if i < 2:
            ax = x + flow_w
            add_textbox(slide, "→",
                        ax, flow_y + Inches(0.12), arrow_w, Inches(0.4),
                        font_name=F_BODY, font_size=22, bold=True,
                        color=C_ACCENT2, align=PP_ALIGN.CENTER)


def _build_metric_card(slide, value, label, left, top, width, height):
    card = add_rect(slide, left, top, width, height,
                    fill_color=RGBColor(0x06, 0x0C, 0x18),
                    line_color=C_ACCENT, line_width_pt=1.0)
    # Big number
    add_textbox(slide, value,
                left + Inches(0.18), top + Inches(0.18),
                width - Inches(0.36), Inches(1.1),
                font_name=F_MONO, font_size=44, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)
    # Label
    add_textbox(slide, label,
                left + Inches(0.18), top + Inches(1.35),
                width - Inches(0.36), Inches(0.55),
                font_name=F_BODY, font_size=15, bold=False,
                color=C_DIM, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# S6 — Investment Opportunity / Closing
# ─────────────────────────────────────────────────────────────────────────────
def build_s6(slide):
    add_bg(slide, BG_ACCENT)
    add_small_logo(slide)

    add_textbox(slide, "05", W - Inches(1.1), MARGIN_T, Inches(0.8), Inches(0.4),
                font_name=F_MONO, font_size=11, color=C_DIM, align=PP_ALIGN.RIGHT)

    add_textbox(slide, "ПОЧЕМУ СЕЙЧАС", MARGIN_L, Inches(0.55), SAFE_W, Inches(0.35),
                font_name=F_MONO, font_size=11, color=C_ACCENT2)

    add_textbox(slide, "Масштабируемая платформа event-рынка",
                MARGIN_L, Inches(0.95), SAFE_W, Inches(0.75),
                font_name=F_HEAD, font_size=34, bold=True, color=C_WHITE)

    add_thin_line(slide, MARGIN_L, Inches(2.0), Inches(0.7))

    # 6 advantage cards 3×2 — no emoji, use accent dots for reliability
    advantages = [
        "Большой и растущий рынок",
        "Простая модель бронирования",
        "AI снижает ручную работу",
        "Объединяет клиентов и подрядчиков",
        "Масштаб по городам и странам",
        "Моб. приложение — App Store / Play",
    ]

    card_w = Inches(3.7)
    card_h = Inches(1.38)
    gap    = Inches(0.20)
    rows   = [Inches(2.25), Inches(2.25) + card_h + gap]
    cols   = [MARGIN_L,
              MARGIN_L + card_w + gap,
              MARGIN_L + 2 * (card_w + gap)]

    for idx, label in enumerate(advantages):
        row = idx // 3
        col = idx % 3
        cx = cols[col]
        cy = rows[row]
        card = add_rect(slide, cx, cy, card_w, card_h,
                        fill_color=RGBColor(0x0A, 0x0F, 0x1C),
                        line_color=RGBColor(0x1A, 0x28, 0x4D), line_width_pt=0.75)
        # accent dot
        add_rect(slide, cx + Inches(0.22), cy + Inches(0.22),
                 Inches(0.08), Inches(0.08),
                 fill_color=C_ACCENT, line_color=None)
        add_textbox(slide, label,
                    cx + Inches(0.22), cy + Inches(0.42),
                    card_w - Inches(0.44), card_h - Inches(0.55),
                    font_name=F_HEAD, font_size=16, bold=True,
                    color=C_WHITE, word_wrap=True)

    # Tagline bar
    bar_y = Inches(5.55)
    bar_h = Inches(0.65)
    bar = add_rect(slide, MARGIN_L, bar_y, SAFE_W, bar_h,
                   fill_color=C_ACCENT, line_color=None)
    add_textbox(slide,
                "Сложная организация мероприятия → простой цифровой сервис",
                MARGIN_L + Inches(0.3), bar_y + Inches(0.1),
                SAFE_W - Inches(0.6), bar_h - Inches(0.15),
                font_name=F_HEAD, font_size=17, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)

    # Investment badge + contact
    badge_w = Inches(3.2)
    badge_h = Inches(0.55)
    badge_x = MARGIN_L
    badge_y = Inches(6.42)
    badge = add_rect(slide, badge_x, badge_y, badge_w, badge_h,
                     fill_color=RGBColor(0x06, 0x0C, 0x18),
                     line_color=C_ACCENT, line_width_pt=1.0)
    add_textbox(slide, "⬡  INVESTMENT OPPORTUNITY",
                badge_x + Inches(0.15), badge_y + Inches(0.1),
                badge_w - Inches(0.3), badge_h - Inches(0.15),
                font_name=F_MONO, font_size=12, bold=True,
                color=C_ACCENT2)

    add_textbox(slide, "+7 747 253 11 38",
                W - Inches(3.8), badge_y + Inches(0.08),
                Inches(3.0), badge_h,
                font_name=F_MONO, font_size=15, bold=True,
                color=C_WHITE, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# Main
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    blank_layout = prs.slide_layouts[6]  # blank layout

    builders = [build_s1, build_s2, build_s3, build_s4, build_s5, build_s6]
    for fn in builders:
        slide = prs.slides.add_slide(blank_layout)
        fn(slide)
        print(f"  OK {fn.__name__}")

    prs.save(OUT_PPTX)
    print(f"\nSaved: {OUT_PPTX}")


if __name__ == "__main__":
    main()
