"""
EVENT AI — Investor Presentation Generator
Generates EVENT-AI-Investor-RU.pptx with 6 slides
Dark premium brand: #04060B bg, #2C7BFF accent, #EEF1F8 text
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy
import os

# ── Constants ─────────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
C_BG       = RGBColor(0x04, 0x06, 0x0B)   # #04060B  main dark bg
C_BG2      = RGBColor(0x0A, 0x11, 0x24)   # #0A1124  secondary bg / gradient-ish
C_ACCENT   = RGBColor(0x2C, 0x7B, 0xFF)   # #2C7BFF  electric blue
C_ACCENT_L = RGBColor(0x5A, 0xA0, 0xFF)   # #5AA0FF  lighter blue
C_TEXT     = RGBColor(0xEE, 0xF1, 0xF8)   # #EEF1F8  almost white
C_SEC      = RGBColor(0x9A, 0xA4, 0xB6)   # #9AA4B6  secondary text
C_MUT      = RGBColor(0x5C, 0x65, 0x7A)   # #5C657A  muted text
C_LINE     = RGBColor(0x22, 0x2D, 0x40)   # hairline dividers
C_CARD     = RGBColor(0x0D, 0x15, 0x27)   # card bg slightly lighter

LOGO_PATH  = r"D:\event-ai\assets\logo-dark.png"
OUT_PATH   = r"D:\event-ai\presentation\EVENT-AI-Investor-RU.pptx"

# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]  # completely blank


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape. Returns the shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()

    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = Pt(line_width if line_width else 0.5)
    else:
        line.fill.background()

    return shape


def set_bg(slide, color=C_BG):
    """Fill the slide background with a solid color rect covering entire slide."""
    shape = slide.shapes.add_shape(
        1, 0, 0, SLIDE_W, SLIDE_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send to back
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)
    return shape


def add_gradient_corner(slide):
    """Add subtle darker rectangle in top-right corner for depth."""
    r = add_rect(slide, 7.5, 0, 5.833, 3.5, fill_color=C_BG2)
    r.line.fill.background()


def add_text_box(slide, text, x, y, w, h,
                 font_name="Montserrat", font_size=16,
                 color=C_TEXT, bold=False, italic=False,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    """Add a text box and return the shape."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return txBox


def add_mono_kicker(slide, text, x, y, w=6):
    """Small mono-font kicker label in accent blue."""
    add_text_box(slide, text, x, y, w, 0.35,
                 font_name="Arial", font_size=11,
                 color=C_ACCENT, bold=True)


def add_title(slide, text, x=0.55, y=0.7, w=12.2, h=1.0, size=40):
    """Large heading in near-white."""
    add_text_box(slide, text, x, y, w, h,
                 font_name="Montserrat", font_size=size,
                 color=C_TEXT, bold=True)


def add_subtitle(slide, text, x=0.55, y=1.55, w=11, h=0.5, size=17, color=C_SEC):
    add_text_box(slide, text, x, y, w, h,
                 font_name="Montserrat", font_size=size,
                 color=color)


def add_hairline(slide, x, y, w, h=0.02):
    """Thin horizontal hairline divider."""
    r = add_rect(slide, x, y, w, h, fill_color=C_LINE)
    r.line.fill.background()


def add_logo_corner(slide, scale=0.7):
    """Small logo top-right corner on non-cover slides."""
    if not os.path.exists(LOGO_PATH):
        return
    logo_w = 1.35 * scale
    logo_h = 0.4  * scale
    slide.shapes.add_picture(
        LOGO_PATH,
        Inches(13.333 - logo_w - 0.25),
        Inches(0.2),
        width=Inches(logo_w),
        height=Inches(logo_h)
    )


def card_with_text(slide, x, y, w, h, text, font_size=15, num_label=None):
    """Rounded card with optional number badge."""
    card = add_rect(slide, x, y, w, h, fill_color=C_CARD, line_color=C_LINE, line_width=0.75)
    # Card text
    tx_x = x + 0.18 + (0.55 if num_label else 0)
    tx_w = w - 0.3  - (0.55 if num_label else 0)
    txBox = slide.shapes.add_textbox(Inches(tx_x), Inches(y + 0.1), Inches(tx_w), Inches(h - 0.15))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "Montserrat"
    run.font.size = Pt(font_size)
    run.font.color.rgb = C_TEXT
    # Number badge
    if num_label:
        nb = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.08), Inches(0.45), Inches(h - 0.15))
        tf2 = nb.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = str(num_label)
        r2.font.name = "Arial"
        r2.font.size = Pt(font_size + 3)
        r2.font.color.rgb = C_ACCENT
        r2.font.bold = True
    return card


def big_number(slide, number_text, label_text, x, y):
    """Large accent number with small label below."""
    add_text_box(slide, number_text, x, y, 3, 0.9,
                 font_name="Arial", font_size=42,
                 color=C_ACCENT, bold=True)
    add_text_box(slide, label_text, x, y + 0.75, 3.5, 0.4,
                 font_name="Montserrat", font_size=13,
                 color=C_SEC)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
def slide1_cover(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)

    # Accent corner rectangle top-right
    r = add_rect(slide, 8.0, 0, 5.333, 7.5, fill_color=C_BG2)
    r.line.fill.background()

    # Vertical accent bar left
    v = add_rect(slide, 0, 0, 0.06, 7.5, fill_color=C_ACCENT)
    v.line.fill.background()

    # Kicker — top
    add_mono_kicker(slide, "ИНВЕСТИЦИОННАЯ ПРЕЗЕНТАЦИЯ", 0.55, 0.45, 7)

    # Logo — large, center-left area
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(
            LOGO_PATH,
            Inches(0.55), Inches(1.15),
            width=Inches(4.5)
        )

    # Subtitle
    add_text_box(
        slide,
        "Умный AI-помощник для организации\nмероприятий под ключ",
        0.55, 3.05, 7.0, 1.1,
        font_name="Montserrat", font_size=20,
        color=C_SEC
    )

    # Hairline divider
    add_hairline(slide, 0.55, 4.35, 6.5)

    # City / year
    add_text_box(slide, "Алматы  ·  Казахстан  ·  2026",
                 0.55, 4.5, 5, 0.4,
                 font_name="Arial", font_size=12,
                 color=C_MUT)

    # Right panel decorative text
    add_text_box(slide, "EVENT AI", 8.4, 2.8, 5, 1.0,
                 font_name="Montserrat", font_size=48,
                 color=RGBColor(0x15, 0x25, 0x42), bold=True)
    add_text_box(slide, "PLATFORM", 8.4, 3.65, 5, 0.5,
                 font_name="Montserrat", font_size=22,
                 color=RGBColor(0x15, 0x25, 0x42), bold=True)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ══════════════════════════════════════════════════════════════════════════════
def slide2_problem(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)
    add_gradient_corner(slide)
    add_logo_corner(slide)

    add_mono_kicker(slide, "ПРОБЛЕМА", 0.55, 0.25)
    add_title(slide,
              "Организация мероприятий сегодня —\nсложный и долгий процесс",
              0.55, 0.55, 11.5, 1.3, size=34)
    add_hairline(slide, 0.55, 1.85, 12.2)

    pains = [
        "Клиент не знает, с чего начать",
        "Площадку, ведущих, артистов, декор, фото-видео нужно искать отдельно",
        "Нет прозрачной сметы",
        "Сложно сравнить подрядчиков",
        "Много времени уходит на звонки и согласования",
    ]

    card_h = 0.74
    gap    = 0.13
    start_y = 2.05

    for i, pain in enumerate(pains):
        y = start_y + i * (card_h + gap)
        # Left accent bar per card
        bar = add_rect(slide, 0.55, y, 0.06, card_h, fill_color=C_ACCENT)
        bar.line.fill.background()
        card = add_rect(slide, 0.68, y, 12.08, card_h, fill_color=C_CARD, line_color=C_LINE, line_width=0.5)
        # Text
        txBox = slide.shapes.add_textbox(Inches(0.95), Inches(y + 0.17), Inches(11.5), Inches(card_h - 0.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = pain
        run.font.name = "Montserrat"
        run.font.size = Pt(16)
        run.font.color.rgb = C_TEXT

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution
# ══════════════════════════════════════════════════════════════════════════════
def slide3_solution(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)
    add_gradient_corner(slide)
    add_logo_corner(slide)

    add_mono_kicker(slide, "РЕШЕНИЕ", 0.55, 0.25)
    add_title(slide, "EVENT AI собирает мероприятие за клиента",
              0.55, 0.55, 11.5, 0.85, size=36)
    add_subtitle(slide,
                 "Пользователь отвечает на простые вопросы в AI-чате —\n"
                 "система автоматически формирует:",
                 0.55, 1.42, 11, 0.7, size=16)
    add_hairline(slide, 0.55, 2.15, 12.2)

    items = [
        "Персональную подборку подрядчиков",
        "Предварительную смету",
        "Сценарий мероприятия",
        "Тайминг",
        "Пригласительные",
        "Бронирование",
    ]

    # 2-column grid: 3 rows x 2 cols
    col1_x = 0.55
    col2_x = 6.75
    card_w = 5.95
    card_h = 1.3
    gap_y  = 0.18
    start_y = 2.35

    for idx, item in enumerate(items):
        col = idx % 2
        row = idx // 2
        x = col1_x if col == 0 else col2_x
        y = start_y + row * (card_h + gap_y)

        # Icon dot
        dot = add_rect(slide, x, y + card_h/2 - 0.07, 0.14, 0.14, fill_color=C_ACCENT)
        dot.line.fill.background()

        card = add_rect(slide, x + 0.25, y, card_w - 0.25, card_h,
                        fill_color=C_CARD, line_color=C_LINE, line_width=0.5)

        txBox = slide.shapes.add_textbox(
            Inches(x + 0.4), Inches(y + 0.28),
            Inches(card_w - 0.55), Inches(card_h - 0.35)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = item
        run.font.name = "Montserrat"
        run.font.size = Pt(17)
        run.font.color.rgb = C_TEXT
        run.font.bold = True

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — How it works
# ══════════════════════════════════════════════════════════════════════════════
def slide4_how(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)
    add_gradient_corner(slide)
    add_logo_corner(slide)

    add_mono_kicker(slide, "КАК ЭТО РАБОТАЕТ", 0.55, 0.25)
    add_title(slide, "От идеи до бронирования —\nв одном приложении",
              0.55, 0.55, 10, 1.0, size=36)
    add_hairline(slide, 0.55, 1.65, 12.2)

    steps = [
        "Пользователь выбирает тип мероприятия",
        "AI задаёт уточняющие вопросы",
        "Система формирует персональную подборку",
        "Клиент выбирает понравившиеся проекты",
        "Платформа помогает забронировать услуги",
    ]

    card_w = 2.15
    card_h = 3.5
    gap_x  = 0.3
    start_x = 0.55
    start_y = 1.9

    for i, step in enumerate(steps):
        x = start_x + i * (card_w + gap_x)

        # Card background
        card = add_rect(slide, x, start_y, card_w, card_h,
                        fill_color=C_CARD, line_color=C_LINE, line_width=0.5)

        # Number circle background
        num_bg = add_rect(slide, x + card_w/2 - 0.32, start_y + 0.22, 0.64, 0.64,
                          fill_color=C_ACCENT)
        num_bg.line.fill.background()

        # Number
        nb = slide.shapes.add_textbox(
            Inches(x + card_w/2 - 0.32), Inches(start_y + 0.22),
            Inches(0.64), Inches(0.64)
        )
        tf2 = nb.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = str(i + 1)
        r2.font.name = "Arial"
        r2.font.size = Pt(20)
        r2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r2.font.bold = True

        # Step text
        txBox = slide.shapes.add_textbox(
            Inches(x + 0.18), Inches(start_y + 1.05),
            Inches(card_w - 0.36), Inches(card_h - 1.2)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = step
        run.font.name = "Montserrat"
        run.font.size = Pt(14)
        run.font.color.rgb = C_TEXT

        # Arrow connector between cards (except last)
        if i < len(steps) - 1:
            arrow_x = x + card_w + 0.05
            arrow = slide.shapes.add_textbox(
                Inches(arrow_x), Inches(start_y + card_h/2 - 0.2),
                Inches(0.22), Inches(0.4)
            )
            tf_a = arrow.text_frame
            pa = tf_a.paragraphs[0]
            pa.alignment = PP_ALIGN.CENTER
            ra = pa.add_run()
            ra.text = ">"
            ra.font.name = "Arial"
            ra.font.size = Pt(18)
            ra.font.color.rgb = C_ACCENT
            ra.font.bold = True

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Monetization
# ══════════════════════════════════════════════════════════════════════════════
def slide5_monetization(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)
    add_gradient_corner(slide)
    add_logo_corner(slide)

    add_mono_kicker(slide, "МОНЕТИЗАЦИЯ", 0.55, 0.25)
    add_title(slide, "Бизнес-модель EVENT AI",
              0.55, 0.55, 10, 0.75, size=38)
    add_hairline(slide, 0.55, 1.38, 12.2)

    # Flow diagram: Клиент → Платформа → Партнёр
    flow_y = 1.55
    for label, x in [("Клиент", 0.5), ("Платформа", 4.8), ("Партнёр", 9.0)]:
        box = add_rect(slide, x, flow_y, 2.6, 0.55,
                       fill_color=C_ACCENT if label == "Платформа" else C_CARD,
                       line_color=C_ACCENT, line_width=0.75)
        tb = slide.shapes.add_textbox(Inches(x), Inches(flow_y), Inches(2.6), Inches(0.55))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Montserrat"
        r.font.size = Pt(15)
        r.font.color.rgb = C_TEXT
        r.font.bold = True

    # Arrows between flow boxes
    for ax in [3.25, 7.55]:
        at = slide.shapes.add_textbox(Inches(ax), Inches(flow_y + 0.08), Inches(0.45), Inches(0.4))
        tf_a = at.text_frame
        pa = tf_a.paragraphs[0]
        pa.alignment = PP_ALIGN.CENTER
        ra = pa.add_run()
        ra.text = "→"
        ra.font.name = "Arial"
        ra.font.size = Pt(20)
        ra.font.color.rgb = C_ACCENT
        ra.font.bold = True

    # Big numbers highlight
    # 7% left
    num7 = slide.shapes.add_textbox(Inches(0.55), Inches(2.35), Inches(2.0), Inches(1.0))
    tf7 = num7.text_frame
    p7 = tf7.paragraphs[0]
    r7 = p7.add_run()
    r7.text = "7%"
    r7.font.name = "Arial"
    r7.font.size = Pt(52)
    r7.font.color.rgb = C_ACCENT
    r7.font.bold = True

    lbl7 = slide.shapes.add_textbox(Inches(0.55), Inches(3.25), Inches(3.2), Inches(0.35))
    tf7l = lbl7.text_frame
    pl = tf7l.paragraphs[0]
    rl = pl.add_run()
    rl.text = "комиссия с подтверждённых заказов"
    rl.font.name = "Montserrat"
    rl.font.size = Pt(12)
    rl.font.color.rgb = C_SEC

    # 50 000 ₸
    num50 = slide.shapes.add_textbox(Inches(3.1), Inches(2.35), Inches(4.0), Inches(1.0))
    tf50 = num50.text_frame
    p50 = tf50.paragraphs[0]
    r50 = p50.add_run()
    r50.text = "50 000 ₸"
    r50.font.name = "Arial"
    r50.font.size = Pt(42)
    r50.font.color.rgb = C_ACCENT
    r50.font.bold = True

    lbl50 = slide.shapes.add_textbox(Inches(3.1), Inches(3.25), Inches(3.5), Inches(0.35))
    tf50l = lbl50.text_frame
    p50l = tf50l.paragraphs[0]
    r50l = p50l.add_run()
    r50l.text = "оплата бронирования"
    r50l.font.name = "Montserrat"
    r50l.font.size = Pt(12)
    r50l.font.color.rgb = C_SEC

    add_hairline(slide, 0.55, 3.65, 12.2)

    # Revenue streams list
    streams = [
        "Платное продвижение партнёров",
        "Премиум-размещение ресторанов, артистов и подрядчиков",
        "B2B-пакеты для event-компаний и площадок",
    ]

    stream_y = 3.85
    for s in streams:
        dot = add_rect(slide, 0.55, stream_y + 0.15, 0.12, 0.12, fill_color=C_ACCENT_L)
        dot.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(stream_y), Inches(11.8), Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = s
        r.font.name = "Montserrat"
        r.font.size = Pt(16)
        r.font.color.rgb = C_TEXT
        stream_y += 0.65

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Why Invest
# ══════════════════════════════════════════════════════════════════════════════
def slide6_why(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    set_bg(slide)
    add_gradient_corner(slide)
    add_logo_corner(slide)

    add_mono_kicker(slide, "ПОЧЕМУ ИНТЕРЕСНО ИНВЕСТОРУ", 0.55, 0.25)
    add_title(slide, "EVENT AI — масштабируемая платформа\nдля event-рынка",
              0.55, 0.55, 11.5, 1.1, size=34)
    add_hairline(slide, 0.55, 1.72, 12.2)

    advantages = [
        "Большой рынок мероприятий",
        "Простая модель бронирования",
        "AI снижает ручную работу",
        "Объединяет клиентов и подрядчиков",
        "Масштабирование по городам и странам",
        "Готовится мобильное приложение для App Store и Play Market",
    ]

    card_w = 3.85
    card_h = 0.95
    gap_x = 0.3
    gap_y = 0.18
    start_x = 0.55
    start_y = 1.92

    for idx, adv in enumerate(advantages):
        col = idx % 3
        row = idx // 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = add_rect(slide, x, y, card_w, card_h,
                        fill_color=C_CARD, line_color=C_LINE, line_width=0.5)

        # Left accent strip
        strip = add_rect(slide, x, y, 0.055, card_h, fill_color=C_ACCENT)
        strip.line.fill.background()

        tb = slide.shapes.add_textbox(
            Inches(x + 0.18), Inches(y + 0.18),
            Inches(card_w - 0.28), Inches(card_h - 0.25)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = adv
        run.font.name = "Montserrat"
        run.font.size = Pt(14)
        run.font.color.rgb = C_TEXT

    # Tagline accent block
    tag_y = 4.22
    tag_bg = add_rect(slide, 0.55, tag_y, 12.2, 1.0, fill_color=C_ACCENT)
    tag_bg.line.fill.background()

    tb_tag = slide.shapes.add_textbox(Inches(0.75), Inches(tag_y + 0.15), Inches(10.0), Inches(0.7))
    tf_tag = tb_tag.text_frame
    tf_tag.word_wrap = True
    p_tag = tf_tag.paragraphs[0]
    p_tag.alignment = PP_ALIGN.CENTER
    r_tag = p_tag.add_run()
    r_tag.text = "EVENT AI превращает сложную организацию мероприятия в простой цифровой сервис"
    r_tag.font.name = "Montserrat"
    r_tag.font.size = Pt(18)
    r_tag.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r_tag.font.bold = True

    # Investment opportunity badge
    badge_y = 5.45
    badge = add_rect(slide, 0.55, badge_y, 3.4, 0.6,
                     fill_color=C_ACCENT, line_color=C_ACCENT, line_width=1.0)
    tb_badge = slide.shapes.add_textbox(Inches(0.55), Inches(badge_y), Inches(3.4), Inches(0.6))
    tf_b = tb_badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.alignment = PP_ALIGN.CENTER
    r_b = p_b.add_run()
    r_b.text = "INVESTMENT OPPORTUNITY"
    r_b.font.name = "Arial"
    r_b.font.size = Pt(13)
    r_b.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r_b.font.bold = True

    # Contact
    add_text_box(slide, "+7 747 253 11 38",
                 4.25, badge_y + 0.1, 4, 0.45,
                 font_name="Arial", font_size=17,
                 color=C_ACCENT, bold=True)

    add_text_box(slide, "eventai.kz  ·  Алматы, Казахстан",
                 8.5, badge_y + 0.1, 4.5, 0.45,
                 font_name="Montserrat", font_size=13,
                 color=C_SEC)

    return slide


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
def build():
    prs = new_prs()

    print("Building slide 1 — Cover...")
    slide1_cover(prs)

    print("Building slide 2 — Problem...")
    slide2_problem(prs)

    print("Building slide 3 — Solution...")
    slide3_solution(prs)

    print("Building slide 4 — How it works...")
    slide4_how(prs)

    print("Building slide 5 — Monetization...")
    slide5_monetization(prs)

    print("Building slide 6 — Why Invest...")
    slide6_why(prs)

    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"\nSaved: {OUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")

    # Verification — re-open and count slides
    prs2 = Presentation(OUT_PATH)
    print(f"Verification — slides in file: {len(prs2.slides)}")

    # Check all text boxes have real text (not images)
    text_count = 0
    for sl in prs2.slides:
        for shape in sl.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            text_count += 1
    print(f"Verification — text runs found: {text_count} (confirms editable text, not images)")


if __name__ == "__main__":
    build()
